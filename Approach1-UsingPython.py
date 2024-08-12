import os
import zipfile
import xml.etree.ElementTree as ET
import pandas as pd
from pptx import Presentation
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from spellchecker import SpellChecker
import re
import textstat
import pandas as pd
from pptx import Presentation
from pptx.util import Pt

def emu_to_pt(emu):
    if emu is None:
        return None
    return emu / 12700


def get_pptx_properties(filepath):
    properties = {}
    with zipfile.ZipFile(filepath, 'r') as pptx_zip:
        with pptx_zip.open('docProps/core.xml') as core_xml:
            tree = ET.parse(core_xml)
            root = tree.getroot()
            ns = {'cp': 'http://schemas.openxmlformats.org/package/2006/metadata/core-properties',
                  'dc': 'http://purl.org/dc/elements/1.1/',
                  'dcterms': 'http://purl.org/dc/terms/'}
            
            properties['title'] = root.find('dc:title', ns).text if root.find('dc:title', ns) is not None else 'N/A'
            properties['author'] = root.find('dc:creator', ns).text if root.find('dc:creator', ns) is not None else 'N/A'
            properties['last_modified_by'] = root.find('cp:lastModifiedBy', ns).text if root.find('cp:lastModifiedBy', ns) is not None else 'N/A'
            properties['created'] = root.find('dcterms:created', ns).text if root.find('dcterms:created', ns) is not None else 'N/A'
            properties['modified'] = root.find('dcterms:modified', ns).text if root.find('dcterms:modified', ns) is not None else 'N/A'
    
    return properties

def list_pptx_files_in_directory(directory):
    pptx_file_list = []

    for filename in os.listdir(directory):
        if filename.endswith('.pptx') and os.path.isfile(os.path.join(directory, filename)):
            pptx_file_list.append(filename)
    
    return pptx_file_list

def print_pptx_files_as_markdown(directory):
    pptx_file_list = list_pptx_files_in_directory(directory)
    pptx_file_count = len(pptx_file_list)

    data = []
    for idx, filename in enumerate(pptx_file_list, start=1):
        filepath = os.path.join(directory, filename)
        properties = get_pptx_properties(filepath)
        data.append([idx, filename, properties['author'], properties['last_modified_by'], properties['created'], properties['modified']])

    df = pd.DataFrame(data, columns=["No.", "File Name", "Author", "Last Modified By", "Created Date", "Modified Date"])

    print(f"The number of .pptx files in the directory '{directory}': {pptx_file_count}\n")
    print(df.to_markdown(index=False))




def print_shape_attributes(shape):
    attributes = dir(shape)
    for attribute in attributes:
        if not attribute.startswith("__"):
            try:
                value = getattr(shape, attribute)
                print(attribute + ":", value)
            except:
                pass

def extract_smartart_text(shape):
    # python-pptx does not support SmartArt directly
    return None, None, ""

def extract_font_and_size(shape):
    paragraphs = shape.text_frame.paragraphs
    font_name = paragraphs[0].font.name if paragraphs[0].font else None
    font_size = paragraphs[0].font.size.pt if paragraphs[0].font.size else None
    text = "\n".join([para.text for para in paragraphs])
    return font_name, font_size, text

def clean_text(text):
    cleaned_text = text.replace('\r', ' ').replace('\n', ' ')
    cleaned_text = re.sub(r'[^a-zA-Z\s]', '', cleaned_text)
    cleaned_text = cleaned_text.lower()
    return cleaned_text

def get_spelling_mistakes(string_list):
    spell = SpellChecker()
    mistakes_dict = {}
    for string in string_list:
        cleaned_string = clean_text(string)
        words = cleaned_string.split()
        for word in words:
            corrected_word = spell.correction(word)
            if corrected_word != word:
                if word not in mistakes_dict:
                    mistakes_dict[word] = corrected_word
    count_mistakes = len(mistakes_dict)
    unique_mistakes = set(mistakes_dict.keys())
    return count_mistakes, mistakes_dict

def calculate_readability(text):
    return textstat.flesch_reading_ease(text)

def get_slide_rank(slide_text):
    clean_text = re.sub(r'\W+', ' ', slide_text)
    readability_score = calculate_readability(clean_text)
    return readability_score

def count_animations(slide):
    # python-pptx does not support animations
    return 0

def count_images(slide):
    return len([shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE])

def count_bullet_points(slide):
    bullet_point_count = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text.startswith("\u2022"):  # Unicode bullet character
                    bullet_point_count += 1
    return bullet_point_count

def is_aligned(shape1, shape2, tolerance=5):
    if abs(shape1.top - shape2.top) <= tolerance or abs((shape1.top + shape1.height) - (shape2.top + shape2.height)) <= tolerance:
        return True
    if abs(shape1.left - shape2.left) <= tolerance or abs((shape1.left + shape1.width) - (shape2.left + shape2.width)) <= tolerance:
        return True
    return False

def count_misalignments_in_slide(slide, tolerance=20):
    misalignments = 0
    shapes = slide.shapes
    for i in range(len(shapes)):
        for j in range(i + 1, len(shapes)):
            shape1 = shapes[i]
            shape2 = shapes[j]
            if not is_aligned(shape1, shape2, tolerance):
                misalignments += 1
    return misalignments


def get_all_placeholder_info(shape):
 
    #placeholder = shape.placeholder_format
    placeholder_text = shape.text
    placeholder_shape = shape.shape_type
    placeholder_size = (shape.width, shape.height)
    font_name = "-"  # Default font name
    font_size = "-"  # Default font size
    if shape.has_text_frame:
        text_frame = shape.text_frame
        if len(text_frame.paragraphs) > 0:
            first_paragraph = text_frame.paragraphs[0]
            if len(first_paragraph.runs) > 0:
                font = first_paragraph.runs[0].font
                if font is not None:
                    font_name = font.name if font.name is not None else font_name
                    font_size = emu_to_pt(font.size)
    return placeholder_text, font_name, font_size
 

def get_first_placeholder_info(slide):
 
    first_slide =slide

    for shape in first_slide.shapes:
        if shape.is_placeholder:
            placeholder = shape.placeholder_format
            if placeholder.idx == 0:  # Assuming you want to get the first placeholder
                placeholder_text = shape.text
                placeholder_shape = shape.shape_type
                placeholder_size = (shape.width, shape.height)
                font_name = "-"  # Default font name
                font_size = "-"  # Default font size
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    if len(text_frame.paragraphs) > 0:
                        first_paragraph = text_frame.paragraphs[0]
                        if len(first_paragraph.runs) > 0:
                            font = first_paragraph.runs[0].font
                            if font is not None:
                                font_name = font.name if font.name is not None else font_name
                                font_size = emu_to_pt(font.size)
                return placeholder_text, font_name, font_size
    return None, None, None

def get_font_info(text_frame):
    fonts = []
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size:
                fonts.append((run.font.name, run.font.size.pt))
    return fonts
def extract_text_and_fonts(shapes):
    text = ''
    fonts = []
    for shape in shapes:
        if hasattr(shape, "text_frame") and shape.text_frame is not None:
            text += shape.text_frame.text + ' '
            fonts.extend(get_font_info(shape.text_frame))
        elif shape.shape_type == 16:  # Shape type 16 corresponds to SmartArt
            text += extract_smartart_text(shape)
    return text, fonts

def read_pptx(file_path):
    presentation = Presentation(file_path)
    results = []
    slide_counter = 1
    for slide in presentation.slides:
        text_list = []
        font_list = []
        size_list = []
        wordcount = 0

        first_shape = None
        for shape in slide.shapes:
            #if shape.has_text_frame and shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            if shape.has_text_frame:
                first_shape = shape
                break
        
        #print(first_shape.text)
        if first_shape:
            headingText,headingfontname,headingsize= get_first_placeholder_info(slide)
            #print(headingfontname, headingsize, headingText)
            text_list.append(headingText)
            try:
                wordcount += len(headingText.split(" "))
            except:
                pass
        else:
            #print("X")
            headingfontname = "NA"
            headingsize = "NA"
            headingText = "NA"

        for shape in slide.shapes:
            if shape.has_text_frame:
                try:
                    placeholder_text, font_name, font_size = get_all_placeholder_info(shape)
                    #print()
                    text_list.append(placeholder_text)
                    font_list.append(font_name)
                    size_list.append(font_size)
                    wordcount += len(placeholder_text.split(" "))
                except:
                    pass
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for shp in shape.shapes:    
                    placeholder_text, font_name, font_size = get_all_placeholder_info(shp)
                    text_list.append(placeholder_text)
                    font_list.append(font_name)
                    size_list.append(font_size)
                    wordcount += len(placeholder_text.split(" "))

        #print(font_list)
        #print(text_list)

        slide_result = {
                        'Slide Number': slide_counter,
                        'Heading Font Name':headingfontname,
                        'Heading Font Size':headingsize,
                        'Text Font Name':set(font_list),
                        'Text Font Size':set(size_list),
                        #'Readability Score':get_slide_rank(" ".join(text_list)),
                        'Word count':wordcount ,
                        'Animation count':count_animations(slide),
                        'Bullet count': count_bullet_points(slide),
                        'Images count':count_images(slide),
                        #'Misalignment':count_misalignments_in_slide(slide, tolerance=5),
                        #'Typos':get_spelling_mistakes(text_list),

                        }
        slide_counter = slide_counter+1
        results.append(slide_result)

    temp = pd.DataFrame(results)


    html_content = temp.to_html()

    # Define output file path for the image
    output_image_path = file_path+'.png'
    import imgkit

    imgkit_options = {
        'format': 'png',
        'width': '1200',  # Width in pixels
        'height': '600',  # Height in pixels
        'quiet': '',  # Suppress wkhtmltoimage command output
    }

    # Use imgkit to convert HTML to image
    imgkit.from_string(html_content, output_image_path,options=imgkit_options)
    imgkit.from_string(html_content, output_image_path)

    print(f"DataFrame converted to image and saved as {output_image_path}")
    print(temp.to_markdown())
    return temp


# Usage example
directory = 'Files'
print_pptx_files_as_markdown(directory)

list_pptx_files_in_directory(directory)

for f in list_pptx_files_in_directory(directory):
    read_pptx(directory+os.sep+f)







