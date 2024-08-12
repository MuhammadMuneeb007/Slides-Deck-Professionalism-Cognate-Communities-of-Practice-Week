import os
import re
import pandas as pd
from collections import Counter
from pptx import Presentation
from spellchecker import SpellChecker
from textblob import TextBlob
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE
import win32com.client

import win32com.client
def print_shape_attributes(shape):
    # Get all attributes of the shape
    attributes = dir(shape)
    
    # Print the name and value of each attribute
    for attribute in attributes:
        # Skip private attributes and methods
        if not attribute.startswith("__"):
            try:
                value = getattr(shape, attribute)
                print(attribute + ":", value)
            except:
                pass
def extract_smartart_text(shape):
    smartart_text = ""
    smartart = shape.SmartArt
    text = []
    for node in smartart.AllNodes:
        #print(node.TextFrame2.TextRange)
        font = node.TextFrame2.TextRange.Font
        #print("Font Name:", font.Name)
        #print("Font Size:", font.Size)
        text.append(node.TextFrame2.TextRange.Text)
        #print(node.TextFrame2.TextRange.Text)
        #try:
        #    smartart_text += node.TextFrame.TextRange.Text + "\n"
        #except:
        #    pass

    return font.Name,font.Size,".".join(text)


def extract_font_and_size(shape):
    text_frame = shape.TextFrame
    text_range = text_frame.TextRange
    font_name = text_range.Font.Name
    font_size = text_range.Font.Size
    return font_name, font_size,text_range.Text

from spellchecker import SpellChecker

import re

def clean_text(text):
    # Remove \r and \n
    cleaned_text = text.replace('\r', ' ').replace('\n', ' ')
    # Remove numbers and special characters
    cleaned_text = re.sub(r'[^a-zA-Z\s]', '', cleaned_text)
    # Convert to lowercase
    cleaned_text = cleaned_text.lower()
    return cleaned_text

def get_spelling_mistakes(string_list):
    # Initialize SpellChecker object
    spell = SpellChecker()
    
    # Dictionary to store mistakes and corrected words
    mistakes_dict = {}
    
    # Iterate through each string in the list
    for string in string_list:
        # Clean the text
        cleaned_string = clean_text(string)
        # Split string into words
        words = cleaned_string.split()
        # Check spelling mistakes for each word
        for word in words:
            # Check if the word is misspelled
            corrected_word = spell.correction(word)
            if corrected_word != word:
                # Update mistakes dictionary
                if word not in mistakes_dict:
                    mistakes_dict[word] = corrected_word
    
    # Get count and unique mistakes
    count_mistakes = len(mistakes_dict)
    unique_mistakes = set(mistakes_dict.keys())
    
    return count_mistakes, mistakes_dict

import re
import textstat


def calculate_readability(text):
    # You can use different readability metrics here
    # For example, Flesch Reading Ease Score or Gunning Fog Index
    return textstat.flesch_reading_ease(text)

def get_slide_rank(slide_text):
    # Remove non-alphanumeric characters
    clean_text = re.sub(r'\W+', ' ', slide_text)

    # Calculate readability score
    readability_score = calculate_readability(clean_text)
    
    # Return the readability score as the rank
    return readability_score


def count_animations(slide):
    return slide.TimeLine.MainSequence.Count

def count_images(slide):
    image_count = 0
    for shape in slide.Shapes:
        if shape.Type == 13:  # Type 13 represents pictures
            image_count += 1
    return image_count


def count_bullet_points(slide):
    bullet_point_count = 0
    for shape in slide.Shapes:
        if shape.HasTextFrame:
            text_frame = shape.TextFrame
            if text_frame.HasText:
                text_range = text_frame.TextRange
                paragraphs_count = text_range.Paragraphs().Count
                for i in range(1, paragraphs_count + 1):
                    #print("XXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXX")
                    paragraph = text_range.Paragraphs(i)
                    #print(paragraph)
                    if paragraph.Text.startswith("\r"):
                        bullet_point_count += 1
    return bullet_point_count
def is_aligned(shape1, shape2, tolerance=5):
    """
    Check if two shapes are aligned within a given tolerance.
    """
    # Check if the top or bottom edges are aligned
    if abs(shape1.Top - shape2.Top) <= tolerance or abs((shape1.Top + shape1.Height) - (shape2.Top + shape2.Height)) <= tolerance:
        return True
    # Check if the left or right edges are aligned
    if abs(shape1.Left - shape2.Left) <= tolerance or abs((shape1.Left + shape1.Width) - (shape2.Left + shape2.Width)) <= tolerance:
        return True
    return False

def count_misalignments_in_slide(slide, tolerance=20):
    """
    Count the number of misalignments in a given slide.
    """
    misalignments = 0
    shapes = slide.Shapes
    shape_count = shapes.Count
    
    for i in range(1, shape_count + 1):
        for j in range(i + 1, shape_count + 1):
            shape1 = shapes.Item(i)
            shape2 = shapes.Item(j)
            if not is_aligned(shape1, shape2, tolerance):
                misalignments += 1
    
    return misalignments
def read_pptx(file_path):
    print(os.path.abspath(file_path))


    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    presentation = powerpoint.Presentations.Open(os.path.abspath(file_path).replace("\\","\\"),WithWindow=0)
    
    #presentation = powerpoint.Presentations.Open("C:\\Users\\kl\\Desktop\\The University of Queensland\\Internship\\Kaplan\\Tuesday - Professional development\\TECH1100 T1 2024-Assessment 03 - Video Submission-Brisbane 4-414509\\"+file_path,WithWindow=0)
    
    results = []
    
    for slide in presentation.Slides:
        #print("Slide", slide.SlideIndex)

        for shape in slide.Shapes:
            #print(shape.Type)
            pass
        text_list = []
        font_list = []
        size_list = []

        for shape in slide.Shapes:

            if shape.HasTextFrame and shape.Type==14:
                first_shape = shape
                break
            else:
                first_shape = None
                
        
        wordcount = 0
 
        if first_shape:
            headingfontname, headingsize,headingText = extract_font_and_size(first_shape)
            text_list.append(headingText)
            wordcount  = wordcount+len(headingText.split(" "))
        else:
            headingfontname ="NA"
            headingsize ="NA"
            headingText ="NA"
            

        
        found_shape_14 = False
        for shape in slide.Shapes:
            if shape.Type == 14 and not found_shape_14:
                found_shape_14 = True
            elif found_shape_14 and shape.HasSmartArt==0:
                try:
                    font_name, font_size, text = extract_font_and_size(shape)
                    text_list.append(text)
                    font_list.append(font_name)
                    size_list.append(font_size)
                    #print(text)
                    wordcount= len(text.split(" "))+wordcount
                except:
                    pass
            elif shape.HasSmartArt==-1:
                # Get font, size, and text
                font_name, font_size, text = extract_smartart_text(shape)
                # Append to lists
                #print(text)
                text_list.append(text)
                font_list.append(font_name)
                size_list.append(font_size)
                wordcount= len(text.split(" "))+wordcount       
 

        #print(text_list)


        slide_result = {
                        'Slide Number': slide.SlideIndex,
                        'Heading Font Name':headingfontname,
                        'Heading Font Size':headingsize,
                        'Text Font Name':set(font_list),
                        'Text Font Size':set(size_list),
                        #'Readability Score':get_slide_rank(" ".join(text_list)),
                        'Word count':wordcount ,
                        'Animation count':count_animations(slide),
                        'Bullet count': count_bullet_points(slide),
                        'Images count':count_images(slide),
                        'Misalignment':count_misalignments_in_slide(slide, tolerance=20),
                        #'Typos':get_spelling_mistakes(text_list),

                        }
        results.append(slide_result)

    temp = pd.DataFrame(results)

    print(temp.to_markdown())
    
    presentation.Close()
    powerpoint.Quit()
    return temp
import matplotlib.pyplot as plt
import pandas as pd
import imgkit
from tabulate import tabulate
import glob
import os

  

a = []
b= []
def list_pptx_files_in_directory(directory):
    pptx_file_list = []

    for filename in os.listdir(directory):
        if filename.endswith('.pptx') and os.path.isfile(os.path.join(directory, filename)):
            pptx_file_list.append(directory+os.sep+filename)
    
    return pptx_file_list


 
root_directory = "Files" 
a = list_pptx_files_in_directory(root_directory)
print(a)


for f in range(0,len(a)):
    print("Processing!")
    #print(a[f].replace("./","")+"//"+b[f])  

    html_content = read_pptx(a[f]).to_html()
    imgkit.from_string(html_content, a[f]+"2.png") 

    continue
    #filename = "BishesGC_234542_assignsubmission_file\\Assessment3BishesGC1814030.pptx"
    #read_pptx(filename)

exit(0)

